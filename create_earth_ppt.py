from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation object
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
DARK_BLUE = RGBColor(25, 51, 102)
LIGHT_BLUE = RGBColor(100, 149, 237)
WHITE = RGBColor(255, 255, 255)
ACCENT = RGBColor(0, 153, 204)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    sp = subtitle_frame.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(32)
    sp.font.color.rgb = LIGHT_BLUE

def add_content_slide(prs, title, content_points):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # Add accent line
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.4), Inches(9), Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.color.rgb = ACCENT
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(8.5), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.level = 0
        p.space_before = Pt(12)

# Slide 1: Title
add_title_slide(prs, "PLANET EARTH", "Exploring Our World")

# Slide 2: Geography
add_content_slide(prs, "Geography", [
    "Surface Area: 510 million km²",
    "Land covers 29% of Earth's surface",
    "Water (oceans, seas, lakes) covers 71%",
    "7 continents and 5 major oceans",
    "Highest point: Mount Everest (8,849 m)",
    "Deepest point: Mariana Trench (10,994 m)"
])

# Slide 3: Climate
add_content_slide(prs, "Climate & Atmosphere", [
    "Atmosphere: 78% Nitrogen, 21% Oxygen, 1% Other",
    "Average surface temperature: 15°C",
    "Climate zones: Tropical, Temperate, Polar",
    "Weather patterns driven by solar energy",
    "Seasons caused by axial tilt",
    "Climate change: Rising temperatures globally"
])

# Slide 4: Natural Features
add_content_slide(prs, "Natural Features", [
    "Massive mountain ranges worldwide",
    "Diverse ecosystems and biomes",
    "Dynamic weather systems",
    "Biodiversity: Millions of species",
    "Freshwater and saltwater systems",
    "Natural wonders: Waterfalls, canyons, forests"
])

# Slide 5: Environment & Life
add_content_slide(prs, "Environment & Life", [
    "Supports over 8 billion humans",
    "Home to millions of plant and animal species",
    "Renewable resources: Water, wind, solar",
    "Environmental challenges: Pollution, deforestation",
    "Conservation efforts worldwide",
    "Sustainable development goals"
])

# Slide 6: Conclusion
add_title_slide(prs, "Our Planet Awaits", "Protect, Preserve, Sustain")

# Save presentation
prs.save('Earth_Presentation.pptx')
print("PowerPoint presentation 'Earth_Presentation.pptx' has been created successfully!")